#!/usr/bin/env python3
"""
Debug why questions 67+ aren't being parsed.
"""
import re
from pptx import Presentation


def extract_question_number(text):
    """Extract question number from text like '1. Question text...'"""
    match = re.match(r'^(\d+)\.\s+', text)
    if match:
        return int(match.group(1))
    return None


def debug_slides(pptx_path, start_slide=130, end_slide=140):
    """Debug specific slides to see their structure."""
    prs = Presentation(pptx_path)
    slides = list(prs.slides)

    for slide_num in range(start_slide, min(end_slide, len(slides))):
        slide = slides[slide_num]
        print(f"\n{'='*80}")
        print(f"SLIDE {slide_num + 1}")
        print(f"{'='*80}")

        for shape_idx, shape in enumerate(slide.shapes):
            if not hasattr(shape, "text_frame"):
                continue

            text = shape.text.strip()[:100]
            print(f"\nShape {shape_idx}: {text}...")

            text_frame = shape.text_frame
            for para_idx, para in enumerate(text_frame.paragraphs):
                para_text = para.text.strip()
                if not para_text:
                    continue

                # Check for bold runs
                bold_count = 0
                for run in para.runs:
                    if run.text.strip() and run.font.bold:
                        bold_count += 1

                if bold_count > 0:
                    print(f"  Para {para_idx}: '{para_text[:60]}' [HAS {bold_count} BOLD RUNS]")
                else:
                    print(f"  Para {para_idx}: '{para_text[:60]}'")


if __name__ == "__main__":
    pptx_path = "question_pptx/Cloud Digital Leader - Practice questions.pptx"

    # Question 66 would be around slide 133 (66*2 + 1)
    # Question 67 would be around slide 135
    # Let's check slides around that area
    print("Checking slides around questions 66-68...")
    debug_slides(pptx_path, 132, 142)
