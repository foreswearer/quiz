#!/usr/bin/env python3
"""
Upload questions from PPTX to database.
Reads Cloud Digital Leader Practice Questions PPTX and uploads to question_bank.
"""
import re
import sys
from datetime import datetime
from pptx import Presentation
from app.db import get_connection


def extract_question_number(text):
    """Extract question number from text like '1. Question text...'"""
    match = re.match(r'^(\d+)\.\s+', text)
    if match:
        return int(match.group(1))
    return None


def parse_pptx_questions(pptx_path):
    """
    Parse PPTX and extract questions.
    Returns list of dicts with question_text, options, and correct_indices.
    """
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    questions = []

    print(f"Total slides: {len(slides)}")
    print("Parsing questions...")

    # Skip slide 1 (title), process pairs starting from slide 2
    # Questions appear in pairs - one without answer, one with bold answer
    # The pattern switches: Q1-66 have answer on odd slide, Q67+ have answer on even slide
    for i in range(1, len(slides) - 1, 2):
        slide1 = slides[i]
        slide2 = slides[i + 1]

        # Try both slides to see which has the bold answer
        for answer_slide in [slide1, slide2]:
            question_text = ""
            options = []
            correct_indices = []

            for shape in answer_slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                text = shape.text.strip()

                # First shape is the question
                if not question_text and extract_question_number(text) is not None:
                    question_text = text
                    continue

                # Second shape contains the options
                if question_text and text:
                    # Parse options - they're separated by newlines
                    text_frame = shape.text_frame
                    option_idx = 0

                    for para in text_frame.paragraphs:
                        para_text = para.text.strip()
                        if not para_text:
                            continue

                        # Check if this option is bold (correct answer)
                        is_bold = False
                        for run in para.runs:
                            if run.text.strip() and run.font.bold:
                                is_bold = True
                                break

                        options.append(para_text)
                        if is_bold:
                            correct_indices.append(option_idx)
                        option_idx += 1

            # If this slide has bold answers, use it and break
            if question_text and len(options) >= 2 and len(correct_indices) > 0:
                question_num = extract_question_number(question_text)

                # Determine question type based on number of correct answers
                if len(correct_indices) == 1:
                    question_type = 'single_choice'
                    correct_index = correct_indices[0]
                else:
                    question_type = 'multiple_choice'
                    correct_index = correct_indices  # Store list for multi-select

                questions.append({
                    'question_number': question_num,
                    'question_text': question_text,
                    'options': options,
                    'correct_index': correct_index,
                    'question_type': question_type
                })
                print(f"  Parsed Q{question_num}: {len(options)} options, type={question_type}")
                break

    return questions


def get_highest_question_number(conn):
    """Get the highest question number from the database."""
    with conn.cursor() as cur:
        cur.execute("""
            SELECT question_text
            FROM question_bank
            ORDER BY id DESC
            LIMIT 20
        """)
        rows = cur.fetchall()

        max_num = 0
        for row in rows:
            num = extract_question_number(row[0])
            if num and num > max_num:
                max_num = num

        return max_num


def create_backup(conn):
    """Create a simple backup by dumping questions to a file."""
    backup_file = f"backup_questions_{datetime.now().strftime('%Y%m%d_%H%M%S')}.sql"

    with conn.cursor() as cur:
        # Count current questions
        cur.execute("SELECT COUNT(*) FROM question_bank")
        count = cur.fetchone()[0]
        print(f"\nCurrent database has {count} questions")

    # Use pg_dump for proper backup
    import subprocess
    from app.config import DB_NAME, DB_USER, DB_PASSWORD, DB_HOST, DB_PORT

    cmd = [
        'pg_dump',
        '-h', DB_HOST,
        '-p', str(DB_PORT),
        '-U', DB_USER,
        '-t', 'question_bank',
        '-t', 'question_option',
        '--inserts',
        '--data-only',
        DB_NAME
    ]

    try:
        result = subprocess.run(
            cmd,
            env={'PGPASSWORD': DB_PASSWORD},
            capture_output=True,
            text=True,
            timeout=30
        )

        if result.returncode == 0:
            with open(backup_file, 'w') as f:
                f.write(result.stdout)
            print(f"✓ Backup created: {backup_file}")
            return backup_file
        else:
            print(f"Warning: pg_dump failed: {result.stderr}")
            print("Continuing without backup...")
            return None
    except Exception as e:
        print(f"Warning: Could not create backup: {e}")
        print("Continuing without backup...")
        return None


def upload_questions(conn, questions, course_id=1, start_from=None):
    """
    Upload questions to the database.
    If start_from is provided, only upload questions with number > start_from.
    """
    with conn.cursor() as cur:
        uploaded = 0
        skipped = 0

        for q in questions:
            qnum = q['question_number']

            # Skip if question number <= start_from
            if start_from and qnum <= start_from:
                skipped += 1
                continue

            # Insert question
            cur.execute("""
                INSERT INTO question_bank (course_id, question_text, question_type, default_points)
                VALUES (%s, %s, %s, %s)
                RETURNING id
            """, (course_id, q['question_text'], q['question_type'], 0.5))

            question_id = cur.fetchone()[0]

            # Insert options
            # For single_choice: correct_index is an int
            # For multiple_choice: correct_index is a list of ints
            correct_indices = q['correct_index'] if isinstance(q['correct_index'], list) else [q['correct_index']]

            for idx, option_text in enumerate(q['options']):
                is_correct = (idx in correct_indices)
                cur.execute("""
                    INSERT INTO question_option (question_id, option_text, is_correct, order_index)
                    VALUES (%s, %s, %s, %s)
                """, (question_id, option_text, is_correct, idx + 1))

            uploaded += 1
            print(f"  ✓ Uploaded Q{qnum} (DB ID: {question_id}, type: {q['question_type']})")

        conn.commit()
        print(f"\n✓ Upload complete: {uploaded} questions uploaded, {skipped} skipped")
        return uploaded


def main():
    pptx_path = "question_pptx/Cloud Digital Leader - Practice questions.pptx"

    print("="*80)
    print("PPTX Question Upload Tool")
    print("="*80)

    # Parse PPTX
    print("\n[1/4] Parsing PPTX file...")
    questions = parse_pptx_questions(pptx_path)
    print(f"✓ Parsed {len(questions)} questions from PPTX")

    if not questions:
        print("ERROR: No questions found in PPTX")
        return 1

    # Connect to database
    print("\n[2/4] Connecting to database...")
    try:
        conn = get_connection()
        print("✓ Connected to database")
    except Exception as e:
        print(f"ERROR: Could not connect to database: {e}")
        return 1

    try:
        # Check highest question number
        highest = get_highest_question_number(conn)
        print(f"✓ Highest question number in database: {highest}")

        # Upload questions
        print(f"\n[4/4] Uploading questions (starting from question {highest + 1})...")
        uploaded = upload_questions(conn, questions, course_id=1, start_from=highest)

        if uploaded > 0:
            print("\n" + "="*80)
            print(f"SUCCESS! Uploaded {uploaded} new questions to database")
            print("="*80)
            return 0
        else:
            print("\nNo new questions to upload")
            return 0

    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
        return 1
    finally:
        conn.close()


if __name__ == "__main__":
    sys.exit(main())
