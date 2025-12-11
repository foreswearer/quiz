"""
Shared PPTX parsing utilities for Cloud Digital Leader quiz questions.
"""
import re
from pptx import Presentation


def extract_question_number(text):
    """Extract question number from text like '1. Question text...'"""
    match = re.match(r"^(\d+)\.\s+", text)
    if match:
        return int(match.group(1))
    return None


def parse_pptx_questions(pptx_path):
    """
    Parse PPTX and extract questions.
    Returns list of dicts with question_text, options, and correct_indices.
    """
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    questions = []

    print(f"Total slides: {len(slides)}")
    print("Parsing questions...")

    # Skip slide 1 (title), process pairs starting from slide 2
    # Questions appear in pairs - one without answer, one with bold answer
    # The pattern switches: Q1-66 have answer on odd slide, Q67+ have answer on even slide
    for i in range(1, len(slides) - 1, 2):
        slide1 = slides[i]
        slide2 = slides[i + 1]

        # Try both slides to see which has the bold answer
        for answer_slide in [slide1, slide2]:
            question_text = ""
            options = []
            correct_indices = []

            for shape in answer_slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                text = shape.text.strip()

                # First shape is the question
                if not question_text and extract_question_number(text) is not None:
                    question_text = text
                    continue

                # Second shape contains the options
                if question_text and text:
                    # Parse options - they're separated by newlines
                    text_frame = shape.text_frame
                    option_idx = 0

                    for para in text_frame.paragraphs:
                        para_text = para.text.strip()
                        if not para_text:
                            continue

                        # Check if this option is bold (correct answer)
                        is_bold = False
                        for run in para.runs:
                            if run.text.strip() and run.font.bold:
                                is_bold = True
                                break

                        options.append(para_text)
                        if is_bold:
                            correct_indices.append(option_idx)
                        option_idx += 1

            # If this slide has bold answers, use it and break
            if question_text and len(options) >= 2 and len(correct_indices) > 0:
                question_num = extract_question_number(question_text)

                # Determine question type based on number of correct answers
                if len(correct_indices) == 1:
                    question_type = "single_choice"
                    correct_index = correct_indices[0]
                else:
                    question_type = "multiple_choice"
                    correct_index = correct_indices  # Store list for multi-select

                questions.append(
                    {
                        "question_number": question_num,
                        "question_text": question_text,
                        "options": options,
                        "correct_index": correct_index,
                        "question_type": question_type,
                    }
                )
                print(
                    f"  Parsed Q{question_num}: {len(options)} options, type={question_type}"
                )
                break

    return questions
