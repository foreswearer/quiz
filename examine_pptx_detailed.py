#!/usr/bin/env python3
"""
Examine PPTX text formatting to find how correct answers are marked.
"""

from pptx import Presentation


def examine_text_formatting(pptx_path):
    """Look for text formatting differences (bold, underline, color)."""
    prs = Presentation(pptx_path)

    # Look at slides 2-5 (questions 1 and 2, each appearing twice)
    for slide_idx in [2, 3, 4, 5]:
        slide = list(prs.slides)[slide_idx - 1]
        print(f"\n{'=' * 80}")
        print(f"SLIDE {slide_idx}")
        print(f"{'=' * 80}")

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                print(
                    f"\nShape: {shape.text[:80] if hasattr(shape, 'text') else 'N/A'}..."
                )
                text_frame = shape.text_frame

                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    for run_idx, run in enumerate(paragraph.runs):
                        if run.text.strip():
                            print(
                                f"  Run {para_idx}.{run_idx}: '{run.text.strip()[:60]}'"
                            )
                            print(f"    Bold: {run.font.bold}")
                            print(f"    Italic: {run.font.italic}")
                            print(f"    Underline: {run.font.underline}")
                            if run.font.color.type is not None:
                                print(f"    Color type: {run.font.color.type}")


if __name__ == "__main__":
    pptx_path = "question_pptx/Cloud Digital Leader - Practice questions.pptx"
    examine_text_formatting(pptx_path)
