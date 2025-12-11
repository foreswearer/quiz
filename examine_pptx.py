#!/usr/bin/env python3
"""
Examine PPTX structure to understand question format.
"""
from pptx import Presentation

def examine_pptx(pptx_path):
    """Examine the structure of the PPTX file."""
    prs = Presentation(pptx_path)

    print(f"Total slides: {len(prs.slides)}")
    print("\nExamining first 3 slides:\n" + "="*80)

    for idx, slide in enumerate(list(prs.slides)[:10], 1):
        print(f"\n--- Slide {idx} ---")
        for shape_idx, shape in enumerate(slide.shapes, 1):
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    print(f"  Shape {shape_idx}:")
                    print(f"    {text[:300]}")
                    if len(text) > 300:
                        print(f"    ... (truncated, total length: {len(text)})")
        print()

if __name__ == "__main__":
    pptx_path = "question_pptx/Cloud Digital Leader - Practice questions.pptx"
    examine_pptx(pptx_path)
